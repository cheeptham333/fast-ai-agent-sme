import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333) # 16:9 widescreen
prs.slide_height = Inches(7.5)

# Fast AI CI Colors
NAVY = RGBColor(11, 27, 61)       # #0B1B3D
CRIMSON = RGBColor(220, 38, 38)   # #DC2626
LIGHT_BG = RGBColor(248, 250, 252)# #F8FAFC
DARK_BG = RGBColor(10, 15, 29)    # #0A0F1D
WHITE = RGBColor(255, 255, 255)
SLATE_GRAY = RGBColor(100, 116, 139) # #64748B
TEXT_DARK = RGBColor(15, 23, 42)  # #0F172A
CARD_BG = RGBColor(241, 245, 249) # #F1F5F9
ACCENT_BLUE = RGBColor(2, 132, 199)# #0284C7
EMERALD = RGBColor(16, 185, 129)  # #10B981

logo_path = "/Users/tri333/Documents/Ebook AI Agent FastAI/public/fast-ai-logo.png"

def add_header(slide, title_text, category_text="FAST AI MASTERCLASS 2026"):
    # Category badge
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = CRIMSON

    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(9.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = NAVY

    # Small Logo on Top Right
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(10.8), Inches(0.4), width=Inches(1.8))

def create_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

# --- SLIDE 1: Cover Slide ---
slide_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(slide_layout)

bg_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = LIGHT_BG
bg_shape.line.fill.background()

strip = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
strip.fill.solid()
strip.fill.fore_color.rgb = CRIMSON
strip.line.fill.background()

if os.path.exists(logo_path):
    slide1.shapes.add_picture(logo_path, Inches(4.8), Inches(1.2), width=Inches(3.7))

title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.7), Inches(10.33), Inches(2.2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "AI Agent & Vibe Coding for SME"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "ขับเคลื่อนด้วย Google Gemini 3.7, Antigravity, ChatGPT 5.6 & Claude Opus 5"
p2.font.size = Pt(18)
p2.font.color.rgb = ACCENT_BLUE
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(10)

card_cov = create_card(slide1, 2.0, 5.2, 9.33, 1.4, bg_color=WHITE, border_color=CARD_BG)
info_box = slide1.shapes.add_textbox(Inches(2.2), Inches(5.35), Inches(8.93), Inches(1.1))
tf_info = info_box.text_frame
p_info = tf_info.paragraphs[0]
p_info.text = "🚀 หลักสูตร Fast AI สำหรับผู้ประกอบการ SME (อายุ 30+ ไม่มีพื้นฐานเขียนโปรแกรม)"
p_info.font.size = Pt(14)
p_info.font.bold = True
p_info.font.color.rgb = NAVY
p_info.alignment = PP_ALIGN.CENTER

p_info2 = tf_info.add_paragraph()
p_info2.text = "KPI: สั่งงานได้จริง • พัฒนาระบบใช้เองได้ใน 1 วัน • ขยายรายได้ ลดรายจ่าย ด้วยต้นทุน 0 บาท"
p_info2.font.size = Pt(12)
p_info2.font.color.rgb = CRIMSON
p_info2.alignment = PP_ALIGN.CENTER
p_info2.space_before = Pt(4)


# --- SLIDE 2: Thailand SME Context 2026 ---
slide2 = prs.slides.add_slide(slide_layout)
add_header(slide2, "บริบท AI ประเทศไทย 2026: โอกาสทองของ SME ยุคใหม่")

create_card(slide2, 0.8, 1.8, 3.7, 4.8, WHITE, ACCENT_BLUE)
tb1 = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.3), Inches(4.4))
tf1 = tb1.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = "📊 ผลสำรวจ AWS 2026"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY

p = tf1.add_paragraph()
p.text = "Unlocking Thailand's AI Potential\n"
p.font.size = Pt(11)
p.font.color.rgb = SLATE_GRAY

p = tf1.add_paragraph()
p.text = "• 84% ของ SME ไทยที่ใช้ AI รายงานว่าประสิทธิภาพงานเพิ่มขึ้นชัดเจน\n\n• 71% มีรายได้เติบโตเฉลี่ย 19% ภายในปีแรกที่นำมาปรับใช้\n\n• ประหยัดเวลาทำงานเอกสาร 15-20 ชม./สัปดาห์"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_DARK

create_card(slide2, 4.8, 1.8, 3.7, 4.8, WHITE, CRIMSON)
tb2 = slide2.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.3), Inches(4.4))
tf2 = tb2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "🇹🇭 UOB Study 2026"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY

p = tf2.add_paragraph()
p.text = "Business Outlook Study 2026\n"
p.font.size = Pt(11)
p.font.color.rgb = SLATE_GRAY

p = tf2.add_paragraph()
p.text = "• ผู้ประกอบการไทยกว่า 70% เริ่มนำ AI เข้าสู่องค์กรแล้ว นำหน้าค่าเฉลี่ยอาเซียน\n\n• SME ที่ไม่ปรับตัวเผชิญปัญหาต้นทุนค่าแรงและตอบลูกค้าช้ากว่าคู่แข่ง 3 เท่า\n\n• การนำ AI Agent มาตอบแชทกู้คืนยอดขายได้ 35%"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_DARK

create_card(slide2, 8.8, 1.8, 3.7, 4.8, NAVY)
tb3 = slide2.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.3), Inches(4.4))
tf3 = tb3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "💡 สิ่งที่ SME ต้องตระหนัก"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = WHITE

p = tf3.add_paragraph()
p.text = "\n1. AI ไม่ได้มาแทนที่คุณ แต่คนที่ใช้ AI เป็นจะมาแทนคนที่ปฏิเสธ AI\n\n2. เครื่องมือยุค 2026 ฟรีและง่าย ไม่จำเป็นต้องรู้โค้ด\n\n3. เริ่มต้นจากโจทย์เล็กๆ ที่เห็นผลกำไรทันที (Quick Win)"
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(226, 232, 240)


# --- SLIDE 3: Vibe Coding Mindset ---
slide3 = prs.slides.add_slide(slide_layout)
add_header(slide3, "Mindset Vibe Coding: เลิกเขียนโค้ด เปลี่ยนเป็น CEO คุม Agent")

create_card(slide3, 0.8, 1.8, 5.6, 5.0, WHITE)
tb_l = slide3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.6))
tf_l = tb_l.text_frame
tf_l.word_wrap = True
p = tf_l.paragraphs[0]
p.text = "🎯 กฎทอง 80/20 ในการสั่งงาน AI Agent"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = CRIMSON

p = tf_l.add_paragraph()
p.text = "\n• 80% ของความสำเร็จ มาจากการ 'ตั้งโจทย์และบรีฟงานที่ชัดเจน' (Context, Target, Business Logic)\n\n• 20% ที่เหลือ คือการตรวจรับงานและขัดเกลา (Feedback Loop)\n\n• Prompt แรกได้ผลงาน 40-60% เสมอ! การปรับแก้ทีละจุด (Iteration) คือหัวใจสำคัญ"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_DARK

create_card(slide3, 6.8, 1.8, 5.7, 5.0, CARD_BG)
tb_r = slide3.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.6))
tf_r = tb_r.text_frame
tf_r.word_wrap = True
p = tf_r.paragraphs[0]
p.text = "🔄 เปรียบเทียบวิธีคิด (Mindset Shift)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY

p = tf_r.add_paragraph()
p.text = "\n❌ แบบเดิม (Developer Mindset):\n- ต้องท่องจำ Syntax, ปีกกา, เซมิโคลอน\n- ใช้เวลาเรียนเป็นปี เสียเวลาธุรกิจ\n\n✅ แบบ Vibe Coding (CEO Mindset):\n- สั่งการด้วยภาษาธุรกิจ (Business Language)\n- AI ทำหน้าที่เป็นทีมวิศวกรซอฟต์แวร์ 5 คน\n- ตรวจงานและอนุมัติด้วยสายตา"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_DARK


# --- SLIDE 4: The Modern Frontier AI Stack (Gemini 3.7, ChatGPT 5.6, Claude Opus 5) ---
slide4 = prs.slides.add_slide(slide_layout)
add_header(slide4, "The Frontier AI Stack: 4 ขุนพลเรือธงประจำองค์กร")

tools_data = [
    ("Google Gemini 3.7", "Primary Engine (พระเอก)", "ขับเคลื่อน Antigravity & NotebookLM ประมวลผลเอกสาร/ภาพ 2M Tokens", NAVY),
    ("Google NotebookLM", "Second Brain (สมองที่ 2)", "คลังความรู้ปลอดข้อมูลมั่ว 100% + Audio Podcast Studio", ACCENT_BLUE),
    ("ChatGPT 5.6 Sol / Canvas", "Strategic Reasoning & UI", "ปรับแต่ง UI สดแบบ Interactive และวิเคราะห์ธุรกิจเชิงลึก", CRIMSON),
    ("Claude Opus 5 / Sonnet 5", "Master Coder & Computer Use", "เขียนโค้ดสถาปัตยกรรมใหญ่ และสั่ง AI ควบคุมหน้าจอคอม", EMERALD),
]

for i, (name, role, desc, col) in enumerate(tools_data):
    left = 0.8 + (i * 2.95)
    create_card(slide4, left, 1.8, 2.8, 4.8, WHITE, col)
    tb = slide4.shapes.add_textbox(Inches(left + 0.15), Inches(2.0), Inches(2.5), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"0{i+1}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = col

    p = tf.add_paragraph()
    p.text = name
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p = tf.add_paragraph()
    p.text = role
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CRIMSON

    p = tf.add_paragraph()
    p.text = f"\n{desc}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK


# --- SLIDE 5: Google Antigravity & Gemini 3.7 Deep Dive ---
slide5 = prs.slides.add_slide(slide_layout)
add_header(slide5, "เจาะลึก Google Antigravity 2.0 & Gemini 3.7 Flash")

create_card(slide5, 0.8, 1.8, 5.7, 5.0, NAVY)
tb = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.6))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "⚡ Slash Commands ทรงพลังประจำองค์กร"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = WHITE

cmds = [
    ("/goal [เป้าหมาย]", "สั่ง Agent ทำงานลึกข้ามคืนจนเสร็จ 100%"),
    ("/schedule", "ตั้งเวลาให้บอทตื่นมาทำงานอัตโนมัติ (Cron Job)"),
    ("/browser", "สั่งเปิดเบราว์เซอร์ส่องเว็บคู่แข่ง ดึงข้อมูล"),
    ("/grill-me", "ให้ AI สัมภาษณ์เราเพื่อตกผลึก Requirement"),
    ("/teamwork-preview", "กระจายงานให้ Subagents หลายตัวทำงานขนานกัน"),
]
for c, d in cmds:
    p = tf.add_paragraph()
    p.text = f"• {c}: {d}"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(226, 232, 240)
    p.space_before = Pt(6)

create_card(slide5, 6.8, 1.8, 5.7, 5.0, WHITE, ACCENT_BLUE)
tb_r = slide5.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.6))
tf_r = tb_r.text_frame
tf_r.word_wrap = True
p = tf_r.paragraphs[0]
p.text = "👥 สถาปัตยกรรมทีมงานย่อย (Subagents)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY

p = tf_r.add_paragraph()
p.text = "\n1. Research Subagent:\nทำหน้าที่สำรวจโครงสร้าง อ่านคู่มือ ส่อง API แบบ Read-only เพื่อวางแผนอย่างรัดกุม\n\n2. Execution Subagent (Self):\nลงมือเขียนไฟล์ สร้างฐานข้อมูล รันคำสั่ง และทดสอบแอปพลิเคชัน\n\n3. Continuous Feedback:\nระบบแจ้งเตือนอัตโนมัติเมื่อเจอปัญหาและทำการแก้โค้ดตัวเอง (Self-Healing)"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_DARK


# --- SLIDE 6: Google NotebookLM ---
slide6 = prs.slides.add_slide(slide_layout)
add_header(slide6, "Google Gemini Notebook (NotebookLM): ศูนย์รวมสมองที่ 2")

create_card(slide6, 0.8, 1.8, 11.733, 5.0, WHITE, ACCENT_BLUE)
tb = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.6))
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "🧠 Second Brain ปลอดการมั่วข้อมูล 100% (Powered by Gemini 3.7)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY

p = tf.add_paragraph()
p.text = "\n1. Source-Grounded Q&A:\nโยนคู่มือบริษัท, แคตตาล็อก, สัญญา, SOP เข้าไป AI จะตอบโดยอ้างอิงเอกสารเท่านั้น พร้อมชี้หน้าและบรรทัดที่มา\n\n2. Audio Overview Studio (Deep Dive Podcast):\nคลิกเดียว แปลงเอกสารหนา 50 หน้า เป็นรายการพอดแคสต์เสียงจำลองภาษาไทย/อังกฤษ เปิดฟังทบทวนขณะขับรถ\n\n3. สะพานเชื่อมสู่ AI Agent:\nสกัด Requirement ทางธุรกิจจาก NotebookLM ส่งต่อให้ Google Antigravity เขียนแอปพลิเคชันได้ตรงเป้า 100%\n\n🔗 ลิงก์คลังความรู้ Fast AI: https://notebook.google.com/notebook/be661c24-57d5-4208-ba66-73e9aff3cca8"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_DARK


# --- SLIDE 7: 5 SME Projects ---
slide7 = prs.slides.add_slide(slide_layout)
add_header(slide7, "5 พิมพ์เขียวโปรเจกต์จริง SME (Workshop Blueprints)")

projects = [
    ("🛍️ บอทส่องราคาคู่แข่ง (Gemini 3.7 Flash)", "Scraper + LINE Notify", "รู้ราคาก่อน ปรับโปรทันที ประหยัด 20 ชม./สัปดาห์"),
    ("🤖 24/7 AI Sales Rep (Gemini 3.7 Flash API)", "Gemini + Sheets CRM", "ตอบแชทกะดึก ปิดการขาย ออเดอร์ไม่ตกหล่น"),
    ("📦 สต็อกอัจฉริยะ (React + Gemini 3.7)", "Smart Inventory + Reorder", "ป้องกันของขาด/ทุนจม คำนวณ Safety Stock"),
    ("📱 Multi-Platform Content Generator", "AI Caption & Graphic", "ผลิต 30 โพสต์ใน 10 นาที ประหยัด 2.5 หมื่น/ด."),
    ("📊 Executive Financial Dashboard", "Finance & Cashflow Chart", "ดูกำไร-ขาดทุน Real-time ไม่ต้องรองบสิ้นเดือน"),
]

for idx, (p_title, p_sub, p_impact) in enumerate(projects):
    top = 1.8 + (idx * 0.98)
    create_card(slide7, 0.8, top, 11.733, 0.85, WHITE)
    tb = slide7.shapes.add_textbox(Inches(1.0), Inches(top + 0.08), Inches(11.3), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"{p_title}  |  {p_sub}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.text = f"🎯 ผลลัพธ์ทางธุรกิจ (ROI): {p_impact}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = CRIMSON


# --- SLIDE 8: Edge Techniques ---
slide8 = prs.slides.add_slide(slide_layout)
add_header(slide8, "เทคนิคระดับว้าว: Computer Use & DeepSeek Local AI")

create_card(slide8, 0.8, 1.8, 5.7, 5.0, WHITE, CRIMSON)
tb1 = slide8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.6))
tf1 = tb1.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = "🖱️ Computer Use (Claude Opus 5)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = CRIMSON

p = tf1.add_paragraph()
p.text = "\n• ให้ AI มองหน้าจอ เลื่อนเมาส์ คลิกปุ่ม และพิมพ์คีย์บอร์ดแทนมนุษย์จริง\n\n• กรอกฟอร์มภาษี e-Filing หรือระบบราชการอัตโนมัติจากไฟล์ Excel\n\n• ล็อกอินโหลดใบเสร็จ/สลิปธนาคารเข้าโฟลเดอร์แยกแผนกทุกสิ้นเดือน\n\n• ส่องหน้าจอคัดลอกข้อมูลข้ามโปรแกรมเก่า (Legacy ERP) เข้า Google Sheets"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_DARK

create_card(slide8, 6.8, 1.8, 5.7, 5.0, WHITE, ACCENT_BLUE)
tb2 = slide8.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.6))
tf2 = tb2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "🐼 DeepSeek-V4 & Local AI (PDPA 100%)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY

p = tf2.add_paragraph()
p.text = "\n• รันโมเดล DeepSeek-V4 แบบ Offline ในคอมพิวเตอร์ออฟฟิศ ข้อมูลไม่ออกสู่คลาวด์\n\n• ประมวลผลข้อมูลเงินเดือน สัญญาความลับ และงบการเงินอย่างปลอดภัย 100%\n\n• Voice-to-Code: ใช้เสียงพูดภาษาไทยบรีฟงานระหว่างเดินทาง ให้ Agent เขียนทั้งระบบเสร็จก่อนถึงออฟฟิศ"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_DARK


# --- SLIDE 9: Production Deployment 0 Baht ---
slide9 = prs.slides.add_slide(slide_layout)
add_header(slide9, "Deploy ขึ้น Production จริง 0 บาท & แผนภาพ Automation")

steps_dep = [
    ("1. Push Code to GitHub", "สั่ง Agent ให้เซฟโค้ดขึ้นตู้เซฟ GitHub อัตโนมัติ"),
    ("2. Connect with Vercel", "กด Add Project ใน Vercel เลือก Repository กด Deploy"),
    ("3. Live in 2 Minutes", "ได้ลิงก์ https://your-app.vercel.app พร้อม SSL ฟรี"),
    ("4. Attach Custom Domain", "ผูกชื่อเว็บของบริษัท เช่น www.mycompany.com ใน 5 นาที"),
]

for i, (st, sd) in enumerate(steps_dep):
    left = 0.8 + (i * 2.95)
    create_card(slide9, left, 1.8, 2.8, 3.2, WHITE)
    tb = slide9.shapes.add_textbox(Inches(left + 0.15), Inches(2.0), Inches(2.5), Inches(2.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Step 0{i+1}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CRIMSON

    p = tf.add_paragraph()
    p.text = st
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p = tf.add_paragraph()
    p.text = f"\n{sd}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK

create_card(slide9, 0.8, 5.3, 11.733, 1.5, CARD_BG, CRIMSON)
tb_sec = slide9.shapes.add_textbox(Inches(1.0), Inches(5.45), Inches(11.3), Inches(1.2))
tf_sec = tb_sec.text_frame
p = tf_sec.paragraphs[0]
p.text = "🔒 กฎเหล็กความปลอดภัย & พ.ร.บ. คุ้มครองข้อมูลส่วนบุคคล (PDPA 2562)"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CRIMSON

p = tf_sec.add_paragraph()
p.text = "• ห้ามใส่ข้อมูลลูกค้าจริง, รหัสผ่าน หรือ API Key ลงใน Prompt เด็ดขาด (ใช้ Mock Data เสมอ)\n• ซ่อน API Key ในไฟล์ .env และใช้ระบบ Authentication ที่ได้มาตรฐานสากล (Firebase / Supabase)"
p.font.size = Pt(11)
p.font.color.rgb = TEXT_DARK


# --- SLIDE 10: 30-Day Action Plan ---
slide10 = prs.slides.add_slide(slide_layout)
add_header(slide10, "แผนปฏิบัติการ 30 วัน (30-Day SME Action Plan)")

weeks = [
    ("สัปดาห์ที่ 1", "เลือก 1 ปัญหาเจ็บปวด (Pain Point) สูงสุดในบริษัทมาทดลองสร้าง MVP"),
    ("สัปดาห์ที่ 2", "ใช้ Google Antigravity & NotebookLM (Gemini 3.7) ร่างระบบและทดสอบ"),
    ("สัปดาห์ที่ 3", "Deploy ขึ้น Vercel ทดลองใช้งานกับลูกค้าหรือพนักงานกลุ่มย่อย"),
    ("สัปดาห์ที่ 4", "วัดผล ROI (เวลาและเงินที่ประหยัดได้) แล้วเริ่มสร้างระบบถัดไป"),
]

for idx, (w_title, w_desc) in enumerate(weeks):
    top = 1.8 + (idx * 1.0)
    create_card(slide10, 0.8, top, 11.733, 0.85, WHITE)
    tb = slide10.shapes.add_textbox(Inches(1.0), Inches(top + 0.1), Inches(11.3), Inches(0.65))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"🗓️ {w_title}: {w_desc}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY

create_card(slide10, 0.8, 6.0, 11.733, 0.85, NAVY)
tb_cta = slide10.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.65))
tf_cta = tb_cta.text_frame
p = tf_cta.paragraphs[0]
p.text = "🚀 เริ่มต้น Vibe Coding วันนี้ — AI ช่วยสร้าง แต่คุณคือผู้นำความสำเร็จสู่องค์กร!"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Save presentation
ppt_output = "/Users/tri333/Documents/Ebook AI Agent FastAI/public/fast_ai_masterclass_presentation.pptx"
prs.save(ppt_output)
print(f"Presentation saved successfully at: {ppt_output}")
